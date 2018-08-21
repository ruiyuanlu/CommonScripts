# coding=utf8
# author: luruiyuan

import os
import re

import numpy as np
import pptx
from pptx.util import Inches, Cm, Pt
from pptx.enum.chart import XL_CHART_TYPE # See https://python-pptx.readthedocs.io/en/latest/api/enum/XlChartType.html#xlcharttype
from pptx.enum.chart import XL_LEGEND_POSITION # See https://python-pptx.readthedocs.io/en/latest/api/enum/XlLegendPosition.html
from pptx.enum.chart import XL_DATA_LABEL_POSITION # See https://python-pptx.readthedocs.io/en/latest/api/enum/XlDataLabelPosition.html
from pptx.chart.data import ChartData, XyChartData, BubbleChartData

class AbstractBasePPT(object):
    '''
    Abstract Base PowerPoint Operator Class for
    different extentions.
    '''

    def _dup_path_check(self, filename):
        '''
        file path duplicate check.
        if duplicate, the later file will be add (number)

        eg: test.pptx -> test(1).pptx -> test(2).pptx

        Args:
            filename: str given file name
        
        Return:
            path: real path of given file after duplication check
        '''
        # calculate path according to the location of logger.py
        par_path, file_name = os.path.split(filename)
        cur_par, _ = os.path.split(__file__)
        dir_path = os.path.join(cur_par, par_path)
        path = os.path.join(dir_path, file_name)
        if not os.path.exists(dir_path): # create dir if neccessary
            os.makedirs(dir_path)
        if os.path.exists(path):
            # 统计重复的文件数目并根据已有文件名构造新的文件名
            prefix, _ , suffix = file_name.rpartition('.')
            path_lis = os.listdir(dir_path)
            regexp = re.compile('%s\((\d+)\).%s' % (prefix, suffix))
            repeat = [int(regexp.findall(p)[0]) for p in path_lis if regexp.match(p)]
            cnt = max(repeat) + 1 if repeat else 1
            dup_filename = '%s(%d).%s' % (prefix, cnt, suffix)
            path = os.path.join(dir_path, dup_filename)

        return os.path.abspath(path)

    def _num2inches(self, *nums):
        '''
        Convert number to inches obejct.

        Args:
            nums: list, tuple all numbers that need to be converted.
        
        Return:
            res: list of converted Inches object.
        '''
        return list(map(lambda n: None if n is None else Inches(float(n)), nums))

    def _num2pts(self, *nums):
        '''
        Convert number to point obejct.
        Point(磅) obj determines font size.

        Args:
            nums: list, tuple all numbers that need to be converted.
        
        Return:
            res: list of converted Pt object.
        '''
        return list(map(lambda n: None if n is None else Pt(float(n)), nums))
    
    def _num2cms(self, *nums):
        '''
        Convert number to Cm obejct.
        Cm(厘米) obj determines length.

        Args:
            nums: list, tuple all numbers that need to be converted.
        
        Return:
            res: list of converted Cm object.
        '''
        return list(map(lambda n: None if n is None else Cm(float(n)), nums))


class SlideLayouts(object):
    '''
    All slide layouts.
    '''
    Title_Slide = 0
    Title_Content = 1
    Section_Header = 2
    Two_Content = 3
    Comparison = 4
    Title_Only = 5
    Blank = 6
    Content_with_Caption = 7
    Picture_with_Caption = 8


class PPTX(AbstractBasePPT):
    '''PPTX operate class

    Bug: Package Not Found Error raised when loading existing empty pptx file.
    '''

    def __init__(self, template=None, file_name='demo.pptx'):
        """
        Init pptx file.
        Or load existing pptx file throung template arg.

        Note: template file cannot be empty pptx

        Args:
            template: None or str Template file path. Or existing pptx file.
            file_name: str Current pptx file name.
        
        Raises:
            Package Not Found: if template arg is an existing empty pptx file.
        """
        super(PPTX, self).__init__()

        self.template = template
        self.filename = super(PPTX, self)._dup_path_check(self._suffix_check(file_name))
        self.prs = pptx.Presentation(self.template) # load existing pptx file

    def _get_slide_layout(self, layout):
        '''Return given layout object.'''
        return self.prs.slide_layouts[layout]
    
    def _convert_nums(self, *nums, unit='cm'):
        '''Return converted values according to unit(长度单位).'''
        convert = self._num2cms if unit.lower() == 'cm' else self._num2inches
        return convert(*nums)
    
    def _suffix_check(self, filename):
        '''Check suffix of filename.'''
        return filename if filename.endswith('.pptx') else ''.join([filename, '.pptx'])

    def save_as(self, save_name):
        '''Save presentation as specified pptx file.
        Note that 'save_as' will overide existing file.
        '''
        self.prs.save(self._suffix_check(save_name))

    def save(self):
        '''Save presentation as file_name attr'''
        self.prs.save(self.filename)

    def add_slide(self, slide_layout):
        '''add new slide. Return the latest slide.

        Args:
            slide_layout: SlideLayout, int The slide layout either be Slidelayout object or int
        '''
        if isinstance(slide_layout, int):
            slide_layout = self._get_slide_layout(slide_layout)
        return self.prs.slides.add_slide(slide_layout)
    
    def add_title_subtitle(self, slide, title, subtitle=None):
        '''
        Add title and subtitle for given slide.
        '''
        if title is not None and hasattr(slide.shapes, 'title'):
            slide.shapes.title.text = str(title) # add title
        if subtitle is not None and len(slide.placeholders) > 2:
            slide.placeholders[1].text = str(subtitle) # add subtitle

    def add_title_slide(self, title='Title', subtitle=None):
        '''
        Add title slide for a pptx.

        Args:
            title: str The title of first slide
            subtitle: str The subtitle of first slide
        
        Return:
            slide: Slide The latest added slide

        See for details:
        https://python-pptx.readthedocs.io/en/latest/user/quickstart.html#hello-world-example
        '''
        slide = self.add_slide(SlideLayouts.Title_Slide)
        self.add_title_subtitle(slide, title, subtitle)
        return slide

    def add_paragraph(self, text_frame, text, level=0, font_size=None, font_bold=False):
        '''
        Add paragraph into given text frame.

        Args:
            text_frame: TextFrame The part of a shape that contains its text.
            text: str the content of paragrph
            level: int 0 <= level <= 8 The level of given paragraph
            font_size: float The size of font of text
            font_bold: bool If bold the font or not
        
        Return:
            p: Paragraph Added paragraph
        '''
        p = text_frame.add_paragraph()
        p.text = text
        p.level = level
        p.font.size = self._num2pts(font_size)[0]
        p.font.bold = font_bold
        return p

    def add_bullet_slide(self, title='Title', text_frame_title='Text frame', paragraphs=None, levels=None):
        '''
        Add bullet_slide.

        Args:
            title: str Title of current bullet slide
            paragraphs: list List of paragraph texts strings, eg: ['text1', 'text2']
            levels: list List of paragraph text level, eg: [1, 2]
        
        Returns:
            slide: slide Bullet slide
            pars: list List of added paragrphas in the slide 
        
        Note: len(paragraphs) == len(levels)

        See for details:
        https://python-pptx.readthedocs.io/en/latest/user/quickstart.html#bullet-slide-example
        '''
        slide = self.add_slide(SlideLayouts.Title_Content)
        shapes = slide.shapes
        shapes.title.text = title # add title
        tf = shapes.placeholders[1].text_frame
        tf.text = text_frame_title # add text frame text
        # add text and set level
        if paragraphs and levels:
            pars = [self.add_paragraph(tf, par, lev) for par, lev in zip(paragraphs, levels)]
        return slide, pars

    def add_picture(self, slide, img_path, left, top, width=None, height=None, unit='cm'):
        '''
        Add a picture into given slide.

        Args:
            slide: Slide object. Given slide to add image.
            img_path: str image path
            left: float Left top location of image.
            top: float Left top location of image.
            width: float Width of resized image.
            height: float Height of resized image.
            unit(长度单位): str Must be 'cm' or 'inches'
        
        Return:
            picture: added resized picture.
        '''
        left, top, width, height = self._convert_nums(left, top, width, height, unit=unit)
        return slide.shapes.add_picture(img_path, left, top, width=width, height=height)

    def add_picture_slide(self, img_path, left, top, width=None, height=None, unit='cm'):
        '''
        Add picture at (left, top), and resize according to (width, height)
        If width and height both are None, use naive image.
        Otherwise, resize image according to given param.

        Note that all calculations based on inches.

        Args:
            img_path: str image path
            left: float Left top location of image.
            top: float Left top location of image.
            width: float Width of resized image.
            height: float Height of resized image.
            unit(长度单位): str Must be 'cm' or 'inches'

        Return:
            slide: the latest added slide containing the picture.
            pic: the added pictue.

        See for details:
        Example:
        https://python-pptx.readthedocs.io/en/latest/user/quickstart.html#add-picture-example
        API:
        https://python-pptx.readthedocs.io/en/latest/api/shapes.html#pptx.shapes.shapetree.GroupShapes.add_picture
        '''
        slide = self.add_slide(SlideLayouts.Blank)
        pic = self.add_picture(slide, img_path, left, top, width=width, height=height, unit=unit)
        return slide, pic

    def add_table(self, slide, rows, cols, left, top, width, height, unit='cm'):
        '''
        Add (rows x cols) table at (left, top), and total width == width, total height == height

        Note that total width == sum of each column width
                  total height == sum of each row height
        
        Args:
            slide: Slide object. Given slide to add table.
            rows: int Number of rows of table
            cols: int Number of columns of table
            left: left top location of table
            top: float Left top location of table
            width: float Total width of table
            height: float Total height of table
            unit(长度单位): str Must be 'cm' or 'inches'
        
        Return:
            table: the table added in the given slide
        
        See for details:
        Example:
        https://python-pptx.readthedocs.io/en/latest/user/quickstart.html#add-table-example
        '''
        left, top, width, height = self._convert_nums(left, top, width, height, unit=unit)
        return slide.shapes.add_table(rows, cols, left, top, width, height).table

    def add_table_slide(self, rows, cols, left, top, width, height, title=None, unit='cm'):
        '''
        Add title only layout slide, and then add (rows x cols) table at (left, top),
        and total width == width, total height == height

        Note that total width == sum of each column width
                  total height == sum of each row height
        
        Args:
            rows: int Number of rows of table
            cols: int Number of columns of table
            left: left top location of table
            top: float Left top location of table
            width: float Total width of table
            height: float Total height of table
            title: str Title of given slide
        
        Returns:
            slide: the latest added slide
            table: the table added in the given slide
        '''
        slide = self.add_slide(SlideLayouts.Title_Only)
        if title:
            slide.shapes.title.text = title
        table = self.add_table(slide, rows, cols, left, top, width, height, unit=unit)
        return slide, table
    
    def fullfill_table(self, table, content):
        """
        Fullfill table using given content matrix.
        Note that content is a matrix that has the same shape as table's.
        If None in content, that cell will be empty.

        Args:
            table: Table Given table to be fufilled
            content: 2D list, ndarray Matrix that contains the element to fufill the table
        
        Return:
            table: Table fufilled table, which is the same object as input table
        """
        c_shape, t_shape = np.shape(content), (len(table.rows), len(table.columns))
        err = "The shape of content '%s' must be the same as the shape of table '%s'."
        assert c_shape == t_shape, err % (str(c_shape), str(t_shape))
        # Fufill table
        for i, j in np.ndindex(c_shape):
            table.cell(i, j).text = str(content[i][j])
        return table

    def add_textbox(self, slide, left, top, width, height, unit='cm'):
        """
        Add textbox to given slide at (left, top) inches location with (width and height).
        The paragraphs can be added into the text frame of added textbox.

        Args:
            slide: Slide object. Given slide to add textbox.
            left: left top location of textbox
            top: float Left top location of textbox
            width: float Total width of textbox
            height: float Total height of textbox
            unit(长度单位): str Must be 'cm' or 'inches'
        
        Return:
            txtbox: Textbox added text box.
        """
        left, top, width, height = self._convert_nums(left, top, width, height, unit=unit)
        return slide.shapes.add_textbox(left, top, width, height)

    def add_textbox_slide(self, left, top, width, height, *paragraphs, unit='cm'):
        '''
        Add a slide containing paragraphs in textbox

        Args:
            left: left top location of textbox
            top: float Left top location of textbox
            width: float Total width of textbox
            height: float Total height of textbox
            paragraphs: list  List of content of each paragraph
            unit(长度单位): str Must be 'cm' or 'inches'

        Returns:
            slide: slide Added slide
            txtbox: Textbox Added text box
            pars: list List of added paragraphs. Empty list if paragraphs not given.
        '''
        slide = self.add_slide(SlideLayouts.Blank)
        txtbox = self.add_textbox(slide, left, top, width, height, unit=unit)
        pars = [self.add_paragraph(txtbox.text_frame, par) for par in paragraphs]
        return slide, txtbox, pars

    def add_legend(self, chart, position='RIGHT', include_in_layout=False):
        '''
        Add legend to given chart.

        Args:
            chart: Chart Given chart
            position: XL_LEGEND_POSITION Position of legend. Default is 'RIGHT'
                All supported are: 'BOTTOM', 'CORNER', 'CUSTOM', 'LEFT', 'RIGHT', 'TOP'
            include_in_layout: bool If legend included in layout. Default is Flase
        
        Return:
            None
        '''
        chart.has_legend = True
        chart.legend.position = getattr(XL_LEGEND_POSITION, position.upper(), XL_LEGEND_POSITION.RIGHT)
        chart.legend.include_in_layout = include_in_layout
    
    def add_data_labels(self, chart, position='BEST_FIT', number_format=None, font_size=None, font_bold=False):
        '''
        Add data labels to given chart.

        Args:
            chart: Chart Given chart
            position: XL_DATA_LABEL_POSITION Position of data labels. Default is 'BEST_FIT'
                All supported are: 'ABOVE', 'BELOW', 'CUSTOM', 'BEST_FIT', 'CENTER', 'INSIDE_BASE'
                'INSIDE_END', 'LEFT', 'MIXED', 'OUTSIDE_END', 'RIGHT'
            number_format: str Format string of numbers. Default is None
            font_size: float,int Font size
            font_bold: bool If bold font
        Return:
            None
        '''
        position = getattr(XL_DATA_LABEL_POSITION, position.upper(), 'BEST_FIT')
        for plt in chart.plots:
            plt.has_data_labels = True
            data_labels = plt.data_labels
            data_labels.position = position
            if number_format is not None:
                data_labels.number_format = str(number_format)
            if font_size is not None:
                data_labels.font.size = self._num2pts(font_size)[0]
            if font_bold is not None:
                data_labels.font.bold = font_bold
    
    def _add_categories_series(self, chart_data, categories, *series):
        '''
        Add categories and series to given chart data

        Args:
            categories: list List of data group names
            series: list List of data series, one of each has 2 element: 'series name', 'data list'
                    eg: *series = *[('Series 1', [1.3, 3.6]),
                                    ('Series 2', [-0.5, 1.6])]
        Return:
            chart_data: ChartData fufilled chart data object
        '''
        chart_data.categories = categories
        for s in series:
            chart_data.add_series(*s)
        return chart_data

    def add_bar_chart(self, slide, x, y, cx, cy, categories, *series, unit='cm'):
        '''
        Add bar chart into given slide.

        Args:
            slide: slide Given slide
            x: x aixs location of chart
            y: float y aixs location of chart
            cx: float Total width of chart
            cy: float Total height of chart
            categories: list List of bar names
            series: list List of data series, one of each has 2 element: 'series name', 'data list'
                    eg: *series = *[('Series 1', [1.3, 3.6]),
                                    ('Series 2', [-0.5, 1.6])]
            unit: str Must 'cm' or 'inches'
        
        Return:
            chart: added chart
        '''
        x, y, cx, cy = self._convert_nums(x, y, cx, cy, unit=unit)
        chart_data = self._add_categories_series(ChartData(), categories, *series)
        graphic_frame = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)
        return graphic_frame.chart
    
    def add_line_chart(self, slide, x, y, cx, cy, categories, *series, unit='cm', smooth=False):
        '''
        Add line chart into given slide.

        Args:
            slide: Slide Given slide
            x: x aixs location of chart
            y: float y aixs location of chart
            cx: float Total width of chart
            cy: float Total height of chart
            categories: list List of data group names
            series: list List of data series, one of each has 2 element: 'series name', 'data list'
                    eg: *series = *[('Series 1', [1.3, 3.6]),
                                    ('Series 2', [-0.5, 1.6])]
            unit: str Must 'cm' or 'inches'
            smooth: bool, bools Bool or list of Bool indicate if smooth lines.
        
        Return:
            line_chart: Line chart
        '''
        x, y, cx, cy = self._convert_nums(x, y, cx, cy, unit=unit)
        chart_data = self._add_categories_series(ChartData(), categories, *series)
        graphic_frame = slide.shapes.add_chart(XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data)
        line_chart = graphic_frame.chart
        # smooth settings
        if isinstance(smooth, bool):
            smooth = [smooth] * len(series)
        else:
            assert len(smooth) == len(series), 'Length of smooth and series not equal.'
        for ser, sm in zip(line_chart.series, smooth):
            ser.smooth = sm
        return line_chart
    
    def _add_series_datapoints(self, chart_data, *series):
        '''
        Add categories and series to given chart data

        Args:
            series: list List of data series, one of each has 2 element: 'series name', 'data list'
                    eg: *series = *[('Series 1', (1.3, 3.6), (2.3, 1.5)),
                                    ('Series 2', (-0.5, 1.6), (0.4, 2.9))]
        Return:
            chart_data: ChartData fufilled chart data object
        '''
        # type() cannot tell diff between inherited classes. Use __class__ instead.
        if chart_data.__class__ == XyChartData:
            for s in series:
                ser = chart_data.add_series(s[0])
                for (x, y) in s[1:]:
                    ser.add_data_point(x, y)
        elif chart_data.__class__ == BubbleChartData:
            for s in series:
                ser = chart_data.add_series(s[0])
                for (x, y, size) in s[1:]:
                    ser.add_data_point(x, y, size)
        else:
            raise NotImplementedError('%s not supported.' % str(type(chart_data)))
        return chart_data

    def add_scatter_chart(self, slide, x, y, cx, cy, *series, unit='cm'):
        '''
        Add scatter chart into given slide.

        Args:
            slide: Slide Given slide
            x: x aixs location of chart
            y: float y aixs location of chart
            cx: float Total width of chart
            cy: float Total height of chart
            series: list List of data series, one of each has 2 element: 'series name', 'point list (x, y)'
                    eg: *series = *[('Series 1', [(1.3, 3.6), (2.3, 1.5)]),
                                    ('Series 2', [(-0.5, 1.6), (0.4, 2.9)])]
            unit: str Must 'cm' or 'inches'
        
        Return:
            scatter_chart: Scatter chart
        '''
        x, y, cx, cy = self._convert_nums(x, y, cx, cy, unit=unit)
        chart_data = self._add_series_datapoints(XyChartData(), *series)
        graphic_frame = slide.shapes.add_chart(XL_CHART_TYPE.XY_SCATTER, x, y, cx, cy, chart_data)
        return graphic_frame.chart
    
    def add_bubble_chart(self, slide, x, y, cx, cy, *series, unit='cm'):
        '''
        Add bubble chart into given slide.

        Args:
            slide: Slide Given slide
            x: x aixs location of chart
            y: float y aixs location of chart
            cx: float Total width of chart
            cy: float Total height of chart
            series: list List of data series, one of each has 2 element: 'series name', 'point list (x, y, size)'
                    eg: *series = *[('Series 1', [(1.3, 3.6, 6.6), (2.3, 1.5, 12)]),
                                    ('Series 2', [(-0.5, 1.6, 15), (0.4, 2.9, 3)])]
            unit: str Must 'cm' or 'inches'
        
        Return:
            bubble_chart: Bubble chart
        '''
        x, y, cx, cy = self._convert_nums(x, y, cx, cy, unit=unit)
        chart_data = self._add_series_datapoints(BubbleChartData(), *series)
        graphic_frame = slide.shapes.add_chart(XL_CHART_TYPE.BUBBLE, x, y, cx, cy, chart_data)
        return graphic_frame.chart
    
    def add_pie_chart(self, slide, x, y, cx, cy, categories, series, unit='cm', legend=True, legend_pos='BOTTOM',
                      include_in_layout=False, label=True, label_pos='OUTSIDE_END', number_format='0.0%'):
        '''
        Add pie chart into given slide.

        Args:
            slide: Slide Given slide
            x: x aixs location of chart
            y: float y aixs location of chart
            cx: float Total width of chart
            cy: float Total height of chart
            categories: list List of data group names
            series: Series Data series that has Only 1 group of data of 2 element: 'series name', 'list of percentages'
                    eg: series = ('Series 1', (0.135, 0.324, 0.180, 0.235, 0.126))
            unit: str Must 'cm' or 'inches'
            legend: bool True is legend is needed, else False. Default is True
            legend_pos: str The position of legend. Default is 'BOTTOM'
            include_in_layout: bool If legend included in the layout. Default is False
            label: bool True is data labels are needed, else False. Default is True
            label_pos: str The position of data labels. Default is 'OUTSIDE_END'
            number_format: str Given number format in the chart data labels
        
        Return:
            pie_chart: Pie chart
        '''
        assert len(series) == 2, 'Only 1 group of data supported.'
        x, y, cx, cy = self._convert_nums(x, y, cx, cy, unit=unit)
        chart_data = self._add_categories_series(ChartData(), categories, series)
        pie_chart = slide.shapes.add_chart(XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data).chart
        if legend:
            self.add_legend(pie_chart, legend_pos, include_in_layout)
        if label:
            self.add_data_labels(pie_chart, label_pos, number_format)
        return pie_chart

    def add_shape(self):
        '''
        This method not implemented yet.
        '''
        raise NotImplementedError('This method not implemented yet.')

    def add_shape_slide(self):
        '''
        This method not implemented yet.
        '''
        raise NotImplementedError('This method not implemented yet.')


if __name__ == '__main__':
    ppt = PPTX(file_name='测试输出的PPT')
    # 添加第一张标题页
    ppt.add_title_slide(title='主标题', subtitle='附属标题')
    # 添加第二张, 演示添加图片页
    cur_slide, pic = ppt.add_picture_slide('./天安门.jpg', left=1, top=2)
    # 为同一页再次添加图片, 这一次对图片进行缩放, 并放置在右下角
    ppt.add_picture(cur_slide, './天安门.jpg', left=20, top=13, width=3, height=5)
    # 添加第三张, 演示分级别说明的bullet页, 级别0-8
    ppt.add_bullet_slide(title='bullet页面大标题', text_frame_title='文字帧的标题',
                         paragraphs=['1级文本', '2级文本', '3级文本', '4级文本\n的下一行'],
                         levels=[0, 1, 2, 3])
    
    # 添加第四张, 演示表格
    slide, table = ppt.add_table_slide(rows=2, cols=2, left=9, top=5, width=6, height=6, title='表格标题')
    # 演示获取表格大小
    print('表格行数: %d 列数 %d' % (len(table.rows), len(table.columns)))
    # 填充表格内容, 表格内容是与表格格式相同的, 传入空字符串 '' 用于将表格内容清空
    content = [['一行一列', '一行二列'], ['二行一列', '二行二列']]
    ppt.fullfill_table(table, content)
    
    # 添加第五张, 演示文本框, 首先添加一个空文本框, 如果对段落格式无要求
    # 可以直接通过 paragraphs 参数添加段落内容, 这里演示单独添加段落内容
    slide, txtbox, _ = ppt.add_textbox_slide(left=2, top=2, width=15, height=12)
    tf = txtbox.text_frame # 获取文字帧
    ppt.add_paragraph(tf, '单独添加的\n1级段落的下一行\n字体20磅并加粗', level=1, font_size=20, font_bold=True)

    # 添加第六张, 图表-柱状图
    slide = ppt.add_slide(SlideLayouts.Title_Only)
    ppt.add_title_subtitle(slide, title='柱状图')
    # 添加柱状图
    bar_chart = ppt.add_bar_chart(slide, 6.5, 3, 14, 12, ['类别1', '类别2'],
                    ('Series 1', [1.4, 2.3]), ('Series 2', [2.1, 1.5]), ('Series 3', [3.5, 4.1]))
    # 为柱状图添加图例
    ppt.add_legend(bar_chart, 'RIGHT')
    # 为柱状图添加数据标签
    ppt.add_data_labels(bar_chart, 'INSIDE_END')
    
    # 添加第七张, 图表-条形图
    slide = ppt.add_slide(SlideLayouts.Title_Only)
    ppt.add_title_subtitle(slide, title='条形图')
    # 添加条形图
    line_chart = ppt.add_line_chart(slide, 6.5, 3, 14, 12,
                    ['一季度', '二季度', '三季度'],
                    ('第一条线 平滑', (32.2, 28.4, 34.7)),
                    ('第二条线 不平滑', (20.4, 18.3, 26.2)), smooth=[True, False])
    # 为条形图添加图例
    ppt.add_legend(line_chart)
    # 添加第八张, 图表-条形图
    slide = ppt.add_slide(SlideLayouts.Title_Only)
    ppt.add_title_subtitle(slide, title='散点图')
    # 添加散点图, 每个类型给2个点, 共给出2组数据
    ppt.add_scatter_chart(slide, 6.5, 3, 14, 12, ('数据一', (32.2, 28.4), (34.7, 12)),
                                                ('数据二', (22.7, 25.4), (31.7, 19.8)))
    
    # 添加第九张, 图表-气泡图
    slide = ppt.add_slide(SlideLayouts.Title_Only)
    ppt.add_title_subtitle(slide, title='气泡图')
    # 添加气泡图, 气泡图与散点图类似, 但是传入的参数是3个
    # 第3个参数指示气泡大小 size
    ppt.add_bubble_chart(slide, 6.5, 3, 14, 12,
                        ('Series 1', (1.3, 3.6, 6.6), (2.3, 1.5, 12)),
                        ('Series 2', (-0.5, 1.6, 15), (0.4, 2.9, 3)))
    
    # 添加第十张, 图表-扇形图
    slide = ppt.add_slide(SlideLayouts.Title_Only)
    ppt.add_title_subtitle(slide, title='扇形图')
    # 添加扇形图, 注意, 扇形图一次只能传入一组数据
    ppt.add_pie_chart(slide, 6.5, 3, 14, 12, ['part 1', 'part 2', 'part 3', 'part 4', 'part 5'],
                      ('Series 1', (0.135, 0.324, 0.180, 0.235, 0.126)))

    # 保存ppt
    ppt.save() # save 方法不会覆盖同名文件, 而是为重复的文件编号, 存储文件名为 '测试输出的PPT'
    ppt.save_as('天安门demo.pptx') # save_as 会覆盖同名文件